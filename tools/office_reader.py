"""OfficeReaderTool — Lee Excel, Word, PDF, PowerPoint y CSV."""
from __future__ import annotations
import logging
from pathlib import Path
from typing import Dict, Any

logger = logging.getLogger(__name__)


class OfficeReaderTool:
    """
    Lee archivos de oficina y los convierte a texto estructurado
    para que los agentes puedan analizarlos.
    """

    SUPPORTED_FORMATS = {".xlsx", ".xls", ".csv", ".docx", ".doc", ".pdf", ".pptx", ".ppt"}

    def read(self, file_path: str) -> Dict[str, Any]:
        """Auto-detecta el formato y lee el archivo."""
        path = Path(file_path)
        suffix = path.suffix.lower()

        if not path.exists():
            raise FileNotFoundError(f"Archivo no encontrado: {file_path}")
        if suffix not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Formato no soportado: {suffix}. Soportados: {self.SUPPORTED_FORMATS}")

        readers = {
            ".xlsx": self._read_excel,
            ".xls": self._read_excel,
            ".csv": self._read_csv,
            ".docx": self._read_word,
            ".doc": self._read_word,
            ".pdf": self._read_pdf,
            ".pptx": self._read_powerpoint,
            ".ppt": self._read_powerpoint,
        }

        reader = readers.get(suffix)
        result = reader(file_path)
        result["file_path"] = file_path
        result["file_type"] = suffix
        logger.info(f"OfficeReader: leído {path.name} ({suffix})")
        return result

    def _read_excel(self, path: str) -> Dict:
        import pandas as pd
        xl = pd.ExcelFile(path)
        sheets = {}
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet_name)
            sheets[sheet_name] = {
                "rows": len(df),
                "columns": list(df.columns),
                "preview": df.head(20).to_markdown(index=False),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
                "stats": df.describe().to_markdown() if len(df.select_dtypes(include="number").columns) > 0 else "",
            }
        return {"sheets": sheets, "sheet_count": len(xl.sheet_names)}

    def _read_csv(self, path: str) -> Dict:
        import pandas as pd
        df = pd.read_csv(path)
        return {
            "rows": len(df),
            "columns": list(df.columns),
            "preview": df.head(20).to_markdown(index=False),
            "stats": df.describe().to_markdown() if len(df.select_dtypes(include="number").columns) > 0 else "",
        }

    def _read_word(self, path: str) -> Dict:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                tables_text.append(" | ".join(cell.text for cell in row.cells))
        return {
            "paragraphs": len(paragraphs),
            "text": "\n".join(paragraphs),
            "tables": "\n".join(tables_text),
        }

    def _read_pdf(self, path: str) -> Dict:
        import pdfplumber
        pages_text = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                pages_text.append(f"[Página {i+1}]\n{text}")
        return {
            "pages": len(pages_text),
            "text": "\n\n".join(pages_text),
        }

    def _read_powerpoint(self, path: str) -> Dict:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            slides.append({"slide": i + 1, "content": "\n".join(slide_text)})
        return {"slides": slides, "slide_count": len(slides)}
